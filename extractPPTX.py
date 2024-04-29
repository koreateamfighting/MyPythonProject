from pptx import Presentation


prs = Presentation(r"C:/Users/MEDIAZEN/Desktop/서울관광명소.pptx")
   # 텍스트 추출 결과 리스트
text_runs = []
# 슬라이드 순회
for idx, slide in enumerate(prs.slides):
    
    # 선택된 슬라이드 내 도형 순회
    for shape in slide.shapes:
        
        # 도형에 텍스트 프레임이 없을 경우
       if not shape.has_text_frame:
            # 다음 도형 선택
            continue

        # 선택된 도형 내 텍스트프레임의 문단 순회
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text_runs.append(run.text)

#ppt_file_path = r"C:/Users/MEDIAZEN/Desktop/서울관광명소.pptx"

print(text_runs)

